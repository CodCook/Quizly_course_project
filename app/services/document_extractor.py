from io import BytesIO
from pathlib import Path

import PyPDF2
from docx import Document
from pptx import Presentation


_SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".pptx", ".ppt"}


def get_file_extension(filename: str) -> str:
    return Path((filename or "").strip()).suffix.lower()


def is_supported_upload(filename: str) -> bool:
    return get_file_extension(filename) in _SUPPORTED_EXTENSIONS


def extract_text_from_pdf(file_bytes: bytes) -> str:
    pdf_file = BytesIO(file_bytes)
    reader = PyPDF2.PdfReader(pdf_file)
    text_parts = []

    for page in reader.pages:
        page_text = page.extract_text() or ""
        if page_text.strip():
            text_parts.append(page_text.strip())

    return "\n\n".join(text_parts).strip()


def extract_text_from_docx(file_bytes: bytes) -> str:
    document = Document(BytesIO(file_bytes))
    lines = [para.text.strip() for para in document.paragraphs if para.text and para.text.strip()]
    return "\n".join(lines).strip()


def extract_text_from_pptx(file_bytes: bytes) -> str:
    presentation = Presentation(BytesIO(file_bytes))
    lines = []

    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text and shape.text.strip():
                lines.append(shape.text.strip())

    return "\n".join(lines).strip()


def extract_text_from_upload(file_bytes: bytes, filename: str) -> str:
    extension = get_file_extension(filename)

    if extension == ".pdf":
        return extract_text_from_pdf(file_bytes)
    if extension == ".docx":
        return extract_text_from_docx(file_bytes)
    if extension == ".pptx":
        return extract_text_from_pptx(file_bytes)
    if extension == ".ppt":
        raise ValueError("Legacy .ppt files are not supported. Please upload .pptx instead")

    raise ValueError("Unsupported file type. Please upload a PDF, DOCX, or PPTX file")
